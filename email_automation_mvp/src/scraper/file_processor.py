import os
import re
import requests
from urllib.parse import urlparse

# Imports for file processing
from pdfminer.high_level import extract_text as extract_pdf_text
import docx
import pptx
import spacy

# Load the spaCy model once when the module is loaded
try:
    NLP = spacy.load('en_core_web_sm')
except OSError:
    print("Warning: Spacy model 'en_core_web_sm' not found. Please run 'python -m spacy download en_core_web_sm'. Name extraction will be skipped.")
    NLP = None


# Regex for finding emails
EMAIL_REGEX = r"[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}"


def download_file(url, save_folder="downloaded_files"):
    """
    Downloads a file from a URL and saves it to a local folder.
    """
    if not os.path.exists(save_folder):
        os.makedirs(save_folder)

    try:
        response = requests.get(url, stream=True, timeout=15)
        response.raise_for_status()  # Raise an exception for bad status codes

        # Get a filename from the URL
        parsed_url = urlparse(url)
        filename = os.path.basename(parsed_url.path)
        if not filename:
            filename = "downloaded_file" + os.path.splitext(parsed_url.path)[1]

        filepath = os.path.join(save_folder, filename)

        with open(filepath, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)

        print(f"  -> Successfully downloaded: {filename}")
        return filepath
    except requests.exceptions.RequestException as e:
        print(f"  -> ERROR: Failed to download file from {url}. Reason: {e}")
        return None

def extract_text_from_file(filepath):
    """
    Extracts text from a file based on its extension.
    Handles .pdf, .docx, and .pptx files.
    """
    if not filepath or not os.path.exists(filepath):
        return None

    try:
        filename = os.path.basename(filepath)
        _, extension = os.path.splitext(filename)
        extension = extension.lower()
        text = ""

        print(f"  -> Extracting text from {filename} (type: {extension})")

        if extension == '.pdf':
            # Limit text size to avoid overwhelming spaCy
            text = extract_pdf_text(filepath, maxpages=50)
        elif extension == '.docx':
            doc = docx.Document(filepath)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            text = '\n'.join(full_text)
        elif extension == '.pptx':
            prs = pptx.Presentation(filepath)
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        full_text.append(shape.text)
            text = '\n'.join(full_text)
        elif extension == '.ppt':
            print(f"  -> WARNING: The legacy .ppt format is not supported. Skipping {filename}.")
            return None
        else:
            print(f"  -> WARNING: Unsupported file type '{extension}'. Skipping {filename}.")
            return None

        # Limit the amount of text passed to spaCy to avoid memory issues
        return text[:100000] # spaCy's limit is 1,000,000 characters

    except Exception as e:
        print(f"  -> ERROR: Failed to extract text from {os.path.basename(filepath)}. Reason: {e}")
        return None
    finally:
        # Clean up the downloaded file after processing
        if os.path.exists(filepath):
            os.remove(filepath)
            print(f"  -> Cleaned up downloaded file: {filename}")


def find_emails_in_text(text):
    """
    Finds all unique email addresses in a block of text.
    """
    if not text:
        return []
    return list(set(re.findall(EMAIL_REGEX, text)))

def find_names_in_text(text):
    """
    Finds potential person names in a block of text using spaCy's NER.
    """
    if not text or not NLP:
        return []

    # Process the text with spaCy
    doc = NLP(text)

    # Extract entities labeled as 'PERSON'
    names = [ent.text.strip() for ent in doc.ents if ent.label_ == 'PERSON']

    # Filter out names that are likely not people
    # e.g., single words, or very long strings
    final_names = []
    for name in names:
        if 1 < len(name.split()) < 4 and len(name) < 30:
            final_names.append(name)

    # Return a list of unique names
    return list(set(final_names))


def process_file_url(url):
    """
    High-level function to orchestrate the downloading, text extraction,
    and information finding from a file URL.
    Wrapped in a try-except block to prevent a single file from crashing the scraper.
    """
    try:
        print(f"  -> Processing document link: {url}")
        filepath = download_file(url)
        if not filepath:
            return [], []

        # The file is deleted inside extract_text_from_file, so get the name for logging now.
        filename_for_logging = os.path.basename(filepath)

        text = extract_text_from_file(filepath)
        if not text:
            return [], []  # File is cleaned up in extract_text_from_file

        emails = find_emails_in_text(text)
        names = find_names_in_text(text)

        if emails:
            print(f"  -> Found {len(emails)} email(s) in {filename_for_logging}")
        if names:
            print(f"  -> Found {len(names)} potential name(s) in {filename_for_logging}")

        return emails, names
    except Exception as e:
        # This is a safety net. More specific errors are handled in the functions below.
        print(f"  -> CRITICAL: An unexpected error occurred while processing file URL {url}. Error: {e}. Skipping file.")
        return [], []
