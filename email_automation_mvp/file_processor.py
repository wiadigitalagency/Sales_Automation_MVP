import os
import re
import requests
from urllib.parse import urlparse

# Imports for file processing
from pdfminer.high_level import extract_text as extract_pdf_text
import docx
import pptx

# Regex for finding emails
EMAIL_REGEX = r"[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}"

# A simple regex for names - looks for capitalized words (e.g., "John Doe")
# This is a basic approach and might not be perfect.
NAME_REGEX = r"([A-Z][a-z]+(?:\s[A-Z][a-z]+)*)"

def download_file(url, save_folder="downloaded_files"):
    """
    Downloads a file from a URL and saves it to a local folder.
    """
    if not os.path.exists(save_folder):
        os.makedirs(save_folder)

    try:
        response = requests.get(url, stream=True, timeout=15)
        response.raise_for_status()  # Raise an exception for bad status codes

        # Get a filename from the URL
        parsed_url = urlparse(url)
        filename = os.path.basename(parsed_url.path)
        if not filename:
            filename = "downloaded_file" + os.path.splitext(parsed_url.path)[1]

        filepath = os.path.join(save_folder, filename)

        with open(filepath, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)

        print(f"  -> Successfully downloaded: {filename}")
        return filepath
    except requests.exceptions.RequestException as e:
        print(f"  -> Failed to download {url}. Error: {e}")
        return None

def extract_text_from_file(filepath):
    """
    Extracts text from a file based on its extension.
    Handles .pdf, .docx, and .pptx files.
    """
    if not filepath or not os.path.exists(filepath):
        return None

    try:
        filename = os.path.basename(filepath)
        _, extension = os.path.splitext(filename)
        extension = extension.lower()
        text = ""

        print(f"  -> Extracting text from {filename} (type: {extension})")

        if extension == '.pdf':
            text = extract_pdf_text(filepath)
        elif extension == '.docx':
            doc = docx.Document(filepath)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            text = '\n'.join(full_text)
        elif extension == '.pptx':
            prs = pptx.Presentation(filepath)
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        full_text.append(shape.text)
            text = '\n'.join(full_text)
        elif extension == '.ppt':
            print(f"  -> WARNING: The legacy .ppt format is not supported. Skipping {filename}.")
            return None
        else:
            print(f"  -> WARNING: Unsupported file type '{extension}'. Skipping {filename}.")
            return None

        return text

    except Exception as e:
        print(f"  -> Failed to extract text from {filename}. Error: {e}")
        return None
    finally:
        # Clean up the downloaded file after processing
        if os.path.exists(filepath):
            os.remove(filepath)
            print(f"  -> Cleaned up downloaded file: {filename}")


def find_emails_in_text(text):
    """
    Finds all unique email addresses in a block of text.
    """
    if not text:
        return []
    return list(set(re.findall(EMAIL_REGEX, text)))

def find_names_in_text(text):
    """
    Finds potential names in a block of text.
    """
    if not text:
        return []
    # This is a simple implementation. A more advanced version could use NLP libraries.
    potential_names = re.findall(NAME_REGEX, text)
    # Filter out very short "names" that are likely just capitalized words
    return [name.strip() for name in potential_names if len(name.strip().split()) > 1]

def process_file_url(url):
    """
    High-level function to orchestrate the downloading, text extraction,
    and information finding from a file URL.
    """
    print(f"  -> Processing document link: {url}")
    filepath = download_file(url)
    if not filepath:
        return [], []

    text = extract_text_from_file(filepath)
    if not text:
        return [], []

    emails = find_emails_in_text(text)
    names = find_names_in_text(text)

    if emails:
        print(f"  -> Found {len(emails)} email(s) in {os.path.basename(filepath)}")
    if names:
        print(f"  -> Found {len(names)} potential name(s) in {os.path.basename(filepath)}")

    return emails, names
